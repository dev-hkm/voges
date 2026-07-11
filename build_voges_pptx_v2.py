# -*- coding: utf-8 -*-
"""
Voges Premium Pitch Deck Generator (python-pptx).
Aligns exactly with the AABW Pitching Playbook 6-slide structure.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# Colors
INK       = RGBColor(0x0F, 0x17, 0x2A)   # Slate 900 (very dark blue-gray text)
PAPER     = RGBColor(0xFF, 0xFF, 0xFF)   # white background
GRAY      = RGBColor(0x64, 0x74, 0x8B)   # Slate 500 (secondary text)
LIGHT     = RGBColor(0xE2, 0xE8, 0xF0)   # Slate 200 (hairlines / dividers)
PANEL     = RGBColor(0xF8, 0xFA, 0xFC)   # Slate 50 (subtle panel fill)
BLACKBG   = RGBColor(0x09, 0x0D, 0x16)   # Deep space black for dark slides
WHITEBG   = RGBColor(0xFF, 0xFF, 0xFF)

# Semantic Colors
SAFE      = RGBColor(0x10, 0xB9, 0x81)   # Emerald 500 (Green)
APPROVE   = RGBColor(0xF5, 0x9E, 0x0B)   # Amber 500 (Amber)
BLOCK     = RGBColor(0xEF, 0x44, 0x44)   # Red 500 (Red)

SAFE_BG   = RGBColor(0xEC, 0xFD, 0xF5)
APPROVE_BG= RGBColor(0xFF, 0xFB, 0xEB)
BLOCK_BG  = RGBColor(0xFE, 0xF2, 0xF2)

# Typography
FONT      = "Segoe UI"
FONT_LT   = "Segoe UI Light"
FONT_SB   = "Segoe UI Semibold"

EMU_W = Inches(13.333)
EMU_H = Inches(7.5)

prs = Presentation()
prs.slide_width = EMU_W
prs.slide_height = EMU_H
BLANK = prs.slide_layouts[6]

# Helpers
def slide():
    return prs.slides.add_slide(BLANK)

def bg(s, color):
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = color

def rect(s, x, y, w, h, fill=None, line=None, line_w=0.75, shape=MSO_SHAPE.RECTANGLE, name=None):
    sp = s.shapes.add_shape(shape, x, y, w, h)
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    if name:
        sp.name = name
    return sp

def txt(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space_after=4, line_spacing=1.05, name=None):
    tb = s.shapes.add_textbox(x, y, w, h)
    if name:
        tb.name = name
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        p.line_spacing = line_spacing
        for (t, size, color, bold, fname) in para:
            r = p.add_run()
            r.text = t
            r.font.size = Pt(size)
            r.font.color.rgb = color
            r.font.bold = bold
            r.font.name = fname
    return tb

def kicker(s, text, color=GRAY):
    txt(s, Inches(0.9), Inches(0.62), Inches(11.5), Inches(0.35),
        [[(text.upper(), 11, color, True, FONT_SB)]], name="Kicker")

def headline(s, text, color=INK, size=28):
    txt(s, Inches(0.9), Inches(0.92), Inches(11.5), Inches(0.85),
        [[(text, size, color, True, FONT)]], name="Headline")

def accent_line(s, color=INK):
    rect(s, Inches(0.9), Inches(0.85), Inches(0.6), Pt(3), fill=color, name="AccentLine")

def footer(s, idx, dark=False):
    tcolor = GRAY if not dark else GRAY
    lcolor = LIGHT if not dark else GRAY
    txt(s, Inches(0.9), Inches(7.02), Inches(8), Inches(0.3),
        [[("VOGES", 8.5, tcolor, True, FONT_SB), ("  ·  AI Agent for GoTyme Bank (Track P3)", 8.5, lcolor, False, FONT)]], name="FooterLeft")
    txt(s, Inches(11.4), Inches(7.02), Inches(1.0), Inches(0.3),
        [[(f"{idx:02d}", 9, tcolor, True, FONT_SB)]], align=PP_ALIGN.RIGHT, name="FooterRight")

def chip(s, x, y, w, label, color, bgc, name=None):
    c = rect(s, x, y, w, Inches(0.34), fill=bgc, line=color, line_w=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE, name=name)
    try: c.adjustments[0] = 0.5
    except Exception: pass
    txt(s, x, y + Inches(0.04), w, Inches(0.28), [[(label, 10.5, color, True, FONT_SB)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

def notes(s, vn, en):
    tf = s.notes_slide.notes_text_frame
    tf.text = "VN — Ghi chú thuyết trình:\n" + vn + "\n\nEN — Speaking script:\n" + en

# ============================================================================
# SLIDE 1 — TEAM + PROMISE (Dark background)
# ============================================================================
s1 = slide(); bg(s1, BLACKBG)
# Accent top lines
rect(s1, Inches(0.9), Inches(0.9), Inches(0.9), Pt(4), fill=SAFE, name="Line1")
rect(s1, Inches(1.9), Inches(0.9), Inches(0.9), Pt(4), fill=APPROVE, name="Line2")
rect(s1, Inches(2.9), Inches(0.9), Inches(0.9), Pt(4), fill=BLOCK, name="Line3")

txt(s1, Inches(0.9), Inches(2.1), Inches(11.5), Inches(1.5),
    [[("Voges", 76, PAPER, True, FONT)]], name="MainTitle")
txt(s1, Inches(0.95), Inches(3.45), Inches(11.5), Inches(0.6),
    [[("Voice-First AI Financial Concierge for GoTyme Bank", 22, RGBColor(0x94, 0xA3, 0xB8), False, FONT_LT)]], name="MainSubtitle")

rect(s1, Inches(0.95), Inches(4.35), Inches(0.6), Pt(3), fill=PAPER, name="TitleDivider")
txt(s1, Inches(0.95), Inches(4.55), Inches(11.5), Inches(0.6),
    [[("Voice is the primary interface. The screen is the safety layer.", 18, PAPER, True, FONT_SB)]], name="TitlePromise")

txt(s1, Inches(0.95), Inches(6.3), Inches(11.5), Inches(0.4),
    [[("GPT Realtime WebRTC  ·  Cloudflare Pages + D1  ·  WebAuthn Biometric Passkeys", 12, GRAY, False, FONT)]], name="TitleStack")
footer(s1, 1, dark=True)
notes(s1,
"Chào ban giám khảo, chúng tôi là đội phát triển Voges. Voges là một trợ lý tài chính bằng giọng nói (Voice-first AI Financial Concierge) dành riêng cho GoTyme Bank. Hãy nhớ thông điệp cốt lõi này: Giọng nói là giao diện tương tác chính, còn màn hình hiển thị đóng vai trò là lớp bảo vệ an toàn.",
"Good afternoon, judges. We are the team behind Voges, a voice-first AI financial concierge built specifically for GoTyme Bank. Our core architectural design is simple: voice is the primary interface, and the screen is the safety layer. We eliminate the friction of complex banking menus while keeping transactions completely secure.")

# ============================================================================
# SLIDE 2 — PROBLEM INSIGHT (Light background)
# ============================================================================
s2 = slide(); bg(s2, PAPER)
kicker(s2, "02 / PROBLEM INSIGHT")
accent_line(s2, color=BLOCK)
headline(s2, "The Friction of Deep Menus & The Risk of Unsecured AI")

# Left Panel (Buried Controls)
rect(s2, Inches(0.9), Inches(2.1), Inches(5.5), Inches(4.2), fill=PANEL, line=LIGHT, line_w=1.0, name="Card_Left")
chip(s2, Inches(1.3), Inches(2.4), Inches(2.5), "BURIED NAVIGATION", BLOCK, BLOCK_BG, name="Chip_Left")
txt(s2, Inches(1.3), Inches(3.0), Inches(4.7), Inches(3.0), [
    [("Buried Controls. ", 16, INK, True, FONT_SB), ("Card settings like freezing cards, enabling online transactions, or adjusting limits require navigating 4+ layers of sub-menus.", 14, GRAY, False, FONT)],
    [("High Friction. ", 16, INK, True, FONT_SB), ("Users in urgent situations (e.g., card lost or fraud suspicion) struggle to locate settings immediately, increasing support costs and user anxiety.", 14, GRAY, False, FONT)]
], space_after=14)

# Right Panel (AI Danger)
rect(s2, Inches(6.9), Inches(2.1), Inches(5.5), Inches(4.2), fill=PANEL, line=LIGHT, line_w=1.0, name="Card_Right")
chip(s2, Inches(7.3), Inches(2.4), Inches(2.5), "THE SECURITY DILEMMA", BLOCK, BLOCK_BG, name="Chip_Right")
txt(s2, Inches(7.3), Inches(3.0), Inches(4.7), Inches(3.0), [
    [("The FAQ Chatbot Trap. ", 16, INK, True, FONT_SB), ("Traditional banking bots only answer static questions. They cannot take action because giving an LLM direct write access is a massive security vulnerability.", 14, GRAY, False, FONT)],
    [("The Hallucination Risk. ", 16, INK, True, FONT_SB), ("If an AI agent can execute transfers or database updates directly based on voice alone, malicious prompts or accidental misinterpretations will result in financial loss.", 14, GRAY, False, FONT)]
], space_after=14)

footer(s2, 2)
notes(s2,
"Phần 2: Nêu rõ nỗi đau (pain point). Ứng dụng ngân hàng hiện nay quá nhiều menu phức tạp, làm giảm trải nghiệm khi cần thao tác nhanh. Ngược lại, nếu đưa AI vào để thực hiện giao dịch tự động bằng giọng nói thì rủi ro bịa đặt (hallucination) hoặc chiếm đoạt tài khoản là cực kỳ cao.",
"Let's look at the problem. Traditional banking apps are feature-rich but suffer from buried navigation. Crucial security settings like freezing a card or toggling online payments are hidden deep in sub-menus. While natural voice control is the obvious solution, financial AI faces a security dilemma: standard chatbots are FAQ-only and useless, while AI agents with direct write access are a massive security liability due to prompt injections and hallucinations.")

# ============================================================================
# SLIDE 3 — AGENTIC WORKFLOW (Light background)
# ============================================================================
s3 = slide(); bg(s3, PAPER)
kicker(s3, "03 / AGENTIC WORKFLOW")
accent_line(s3, color=SAFE)
headline(s3, "Zero-Trust Execution: Propose on Client, Decide on Server")

steps = [
    ("1. GOAL", "Speech Stream", "User speaks naturally over direct WebRTC audio connection."),
    ("2. PLAN", "Intent Parsing", "GPT Realtime detects intent & proposes a structured tool call."),
    ("3. ROUTE", "Tool Router", "Server router intercepts write proposal, bypassing client control."),
    ("4. ACT", "Pending Action", "Server creates expiring pending action and displays approval sheet."),
    ("5. VERIFY", "WebAuthn & D1", "User signs with hardware passkey. Server verifies and executes.")
]

bw = Inches(2.15); bh = Inches(3.4)
x0 = Inches(0.9); y0 = Inches(2.2); gap = Inches(0.2)

for i, (title, sub, desc) in enumerate(steps):
    x = x0 + i * (bw + gap)
    rect(s3, x, y0, bw, bh, fill=PANEL, line=LIGHT, line_w=1.0, name=f"Flow_{i+1}")
    rect(s3, x, y0, bw, Inches(0.08), fill=SAFE)
    txt(s3, x + Inches(0.15), y0 + Inches(0.25), bw - Inches(0.3), Inches(0.4),
        [[(title, 14, SAFE, True, FONT_SB)]])
    txt(s3, x + Inches(0.15), y0 + Inches(0.65), bw - Inches(0.3), Inches(0.4),
        [[(sub, 14.5, INK, True, FONT_SB)]])
    txt(s3, x + Inches(0.15), y0 + Inches(1.15), bw - Inches(0.3), Inches(2.1),
        [[(desc, 12, GRAY, False, FONT)]], line_spacing=1.1)

# Draw connecting arrows between boxes
for i in range(len(steps) - 1):
    ax = x0 + (i + 1) * bw + i * gap
    ay = y0 + bh/2
    ln = s3.shapes.add_connector(2, 0, 0, 0, 0)
    ln.begin_x, ln.begin_y, ln.end_x, ln.end_y = int(ax), int(ay), int(ax + gap), int(ay)
    ln.line.color.rgb = GRAY
    ln.line.width = Pt(1.5)
    tail = ln.line._get_or_add_ln()
    end = tail.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'})
    tail.append(end)

footer(s3, 3)
notes(s3,
"Phần 3: Kiến trúc Agent của Voges. Luồng xử lý chia làm 5 bước rõ ràng: Giọng nói truyền qua WebRTC -> GPT đề xuất Tool Call -> Server Router chặn lại và chuyển tới Policy Engine -> Tạo Pending Action trên màn hình -> Người dùng chạm Passkey để xác nhận và thực thi trên D1.",
"This slide shows the Zero-Trust agentic workflow we built. It is a 5-step pipeline. First, the user states their Goal over WebRTC. Second, GPT Realtime plans and proposes a tool call. Third, our backend Tool Router intercepts it. Fourth, if it is sensitive, we act by generating an expiring Pending Action on screen. Fifth, the user verifies using a hardware passkey, and the server executes the SQL transaction in Cloudflare D1. The LLM only proposes; it never directly executes.")

# ============================================================================
# SLIDE 4 — WHY IT WINS (Dark background for contrast)
# ============================================================================
s4 = slide(); bg(s4, BLACKBG)
kicker(s4, "04 / WHY IT WINS")
accent_line(s4, color=APPROVE)
headline(s4, "Three Pillars of Zero-Trust AI Banking", color=PAPER)

pillars = [
    ("POLICY-CONTROLLED", SAFE, SAFE_BG,
     "Deterministic Rules",
     "A server-side, code-authoritative engine evaluates all actions. LLM hallucinations are physically blocked. High-risk writes like money transfers, bypass attempts, and KYC changes are rejected outright by deterministic rules."),
    
    ("PASSKEY-SECURED", APPROVE, APPROVE_BG,
     "Cryptographic Identity",
     "Native WebAuthn integration binds sessions. The server only receives and verifies cryptographic signatures. No raw fingerprint or face data ever leaves the user's secure hardware enclave, preserving biometric privacy."),
    
    ("AUDIT-BACKED", BLOCK, BLOCK_BG,
     "Immutable Logs",
     "Every speech turn, tool call, policy check, WebAuthn assertion, and database transaction outcome is recorded in an append-only, tamper-proof audit log. Perfect for compliance, debugging, and post-action review.")
]

pw = Inches(3.64); ph = Inches(4.2)
px0 = Inches(0.9); py0 = Inches(2.1); pgap = Inches(0.3)

for i, (tag, col, bgc, title, desc) in enumerate(pillars):
    x = px0 + i * (pw + pgap)
    # Background panel
    rect(s4, x, py0, pw, ph, fill=BLACKBG, line=col, line_w=1.2, shape=MSO_SHAPE.ROUNDED_RECTANGLE, name=f"Pillar_{i+1}")
    chip(s4, x + Inches(0.35), py0 + Inches(0.35), Inches(2.2), tag, col, bgc, name=f"PillarChip_{i+1}")
    txt(s4, x + Inches(0.35), py0 + Inches(0.95), pw - Inches(0.7), Inches(0.5),
        [[(title, 18, PAPER, True, FONT_SB)]])
    txt(s4, x + Inches(0.35), py0 + Inches(1.5), pw - Inches(0.7), Inches(2.5),
        [[(desc, 13.5, RGBColor(0x94, 0xA3, 0xB8), False, FONT)]], line_spacing=1.12)

footer(s4, 4, dark=True)
notes(s4,
"Phần 4: Ba trụ cột cốt lõi làm nên sức mạnh của Voges: (1) Policy-Controlled - kiểm soát chặt chẽ phía máy chủ; (2) Passkey-Secured - sử dụng xác thực sinh trắc học phần cứng an toàn tuyệt đối; (3) Audit-Backed - nhật ký hệ thống không thể sửa xóa giúp tăng cường tính minh bạch.",
"Voges wins because of three strict design pillars. First, it is Policy-Controlled: deterministic backend rules override the LLM. Hallucinations cannot compromise security. Second, it is Passkey-Secured: we verify identity cryptographically using WebAuthn. Fingerprints and face data never touch our servers. Third, it is Audit-Backed: an append-only audit trail records every intent, block, and state change for total auditability.")

# ============================================================================
# SLIDE 5 — EVIDENCE + IMPACT (Light background)
# ============================================================================
s5 = slide(); bg(s5, PAPER)
kicker(s5, "05 / EVIDENCE + IMPACT")
accent_line(s5, color=SAFE)
headline(s5, "Rigorous Safety Benchmarks & GoTyme Product Value")

# Left: Safety Benchmarks
rect(s5, Inches(0.9), Inches(2.1), Inches(5.5), Inches(4.2), fill=PANEL, line=LIGHT, line_w=1.0, name="Evidence_Left")
chip(s5, Inches(1.3), Inches(2.4), Inches(2.5), "SECURITY BENCHMARKING", SAFE, SAFE_BG)
txt(s5, Inches(1.3), Inches(2.95), Inches(4.7), Inches(3.0), [
    [("34 / 34 Automated Tests Pass. ", 15.5, INK, True, FONT_SB), ("We run an extensive automated suite validating rules against frozen accounts, unauthorized limits, high-risk transfers, and bypass words.", 13.5, GRAY, False, FONT)],
    [("Zero Leakage of Secrets. ", 15.5, INK, True, FONT_SB), ("100% success rate in deterministic filters blocking the leakage of CVV, OTP, and KYC profile details.", 13.5, GRAY, False, FONT)],
    [("Scam Risk Advisor. ", 15.5, INK, True, FONT_SB), ("Built-in RAG-like check matches conversation context with scam databases to proactively advise users on potential fraud attempts.", 13.5, GRAY, False, FONT)]
], space_after=10)

# Right: Business Value
rect(s5, Inches(6.9), Inches(2.1), Inches(5.5), Inches(4.2), fill=PANEL, line=LIGHT, line_w=1.0, name="Evidence_Right")
chip(s5, Inches(7.3), Inches(2.4), Inches(2.5), "GOTYME PRODUCT UTILITY", SAFE, SAFE_BG)
txt(s5, Inches(7.3), Inches(2.95), Inches(4.7), Inches(3.0), [
    [("90% Navigation Savings. ", 15.5, INK, True, FONT_SB), ("Complex security settings (like turning on online payments) take 5 seconds of natural voice conversation compared to 1 minute of menu searching.", 13.5, GRAY, False, FONT)],
    [("Seeded Banking Scenarios. ", 15.5, INK, True, FONT_SB), ("Integrated with GoTyme Everyday Account, Go Save (5% interest high-yield goals), and PAX Gold (tokenized physical gold) for real product context.", 13.5, GRAY, False, FONT)],
    [("Instant Card Printing. ", 15.5, INK, True, FONT_SB), ("If a card is frozen and replaced, the app immediately generates a support ticket for physical printing at any GoTyme retail kiosk.", 13.5, GRAY, False, FONT)]
], space_after=10)

footer(s5, 5)
notes(s5,
"Phần 5: Các chỉ số thực tế và tính thực tiễn. Đã test thành công 34 kịch bản tự động bao gồm bảo mật, phát hiện lừa đảo (Scam Advisor). Giảm 90% thời gian thao tác. Tích hợp sâu vào sản phẩm GoTyme như tài khoản tiết kiệm Go Save 5% lãi suất và tài khoản Gold PAXG.",
"To validate credibility, we don't just rely on slides. We built a test suite with 34 automated scenarios covering scam risk, frozen account checks, and international payments. All 34 tests pass successfully. In terms of impact, Voges reduces settings navigation time by 90%. We also seeded our database with GoTyme's actual product lines, including Everyday Accounts, high-yield Go Save accounts, and tokenized PAX Gold accounts.")

# ============================================================================
# SLIDE 6 — DEMO + CLOSE (Dark background)
# ============================================================================
s6 = slide(); bg(s6, BLACKBG)
kicker(s6, "06 / DEMO + CLOSE")
accent_line(s6, color=APPROVE)
headline(s6, "Live 3-Step Demo & Production Roadmap", color=PAPER)

# Left Column (Demo paths)
rect(s6, Inches(0.9), Inches(2.1), Inches(5.5), Inches(4.2), fill=BLACKBG, line=GRAY, line_w=1.0, name="Demo_Left")
txt(s6, Inches(1.25), Inches(2.35), Inches(4.8), Inches(0.4),
    [[("LIVE DEMO WALKTHROUGH", 12, RGBColor(0x94, 0xA3, 0xB8), True, FONT_SB)]])
txt(s6, Inches(1.25), Inches(2.85), Inches(4.8), Inches(3.2), [
    [("1. READ (Green)  ·  ", 14.5, SAFE, True, FONT_SB), ("“Why was my Netflix payment declined?”", 14.5, PAPER, True, FONT_SB)],
    [("AI reads recent transactions, diagnoses that online payments are disabled, and explains it naturally.", 12.5, RGBColor(0x94, 0xA3, 0xB8), False, FONT)],
    
    [("2. WRITE (Amber)  ·  ", 14.5, APPROVE, True, FONT_SB), ("“Enable online payments.”", 14.5, PAPER, True, FONT_SB)],
    [("AI proposes tool. Policy approves confirmation. Approval sheet appears. User touches biometric Passkey to execute.", 12.5, RGBColor(0x94, 0xA3, 0xB8), False, FONT)],
    
    [("3. BLOCK (Red)  ·  ", 14.5, BLOCK, True, FONT_SB), ("“Transfer all money without limits.”", 14.5, PAPER, True, FONT_SB)],
    [("Policy engine immediately intercepts, blocks action, logs the threat, and AI refuses safely.", 12.5, RGBColor(0x94, 0xA3, 0xB8), False, FONT)]
], space_after=6)

# Right Column (Roadmap)
rect(s6, Inches(6.9), Inches(2.1), Inches(5.5), Inches(4.2), fill=BLACKBG, line=GRAY, line_w=1.0, name="Demo_Right")
txt(s6, Inches(7.25), Inches(2.35), Inches(4.8), Inches(0.4),
    [[("ROADMAP TO PRODUCTION", 12, RGBColor(0x94, 0xA3, 0xB8), True, FONT_SB)]])
txt(s6, Inches(7.25), Inches(2.85), Inches(4.8), Inches(3.2), [
    [("1. Identity Binding. ", 15.5, PAPER, True, FONT_SB), ("Transition from a single hardcoded demo customer profile to dynamic customer authentication during WebRTC handshake.", 13.5, RGBColor(0x94, 0xA3, 0xB8), False, FONT)],
    [("2. Multi-Customer Sessions. ", 15.5, PAPER, True, FONT_SB), ("Add support for secure session switching and multi-account authorization boundaries once identity is bound.", 13.5, RGBColor(0x94, 0xA3, 0xB8), False, FONT)],
    [("3. Core Banking Integration. ", 15.5, PAPER, True, FONT_SB), ("Connect tool executors to real GoTyme ledgers, card hosts, and support ticket queues via secure API gateways.", 13.5, RGBColor(0x94, 0xA3, 0xB8), False, FONT)]
], space_after=10)

footer(s6, 6, dark=True)
notes(s6,
"Phần 6: Kịch bản demo trực tiếp và Lộ trình phát triển. Demo gồm 3 bước: Hỏi tại sao Netflix decline (Read), Bật online payments (Write - cần xác nhận và passkey), Chuyển khoản hoặc bypass bảo mật (Block - bị chặn). Lộ trình đưa lên sản xuất gồm liên kết Identity thật, hỗ trợ nhiều tài khoản, và kết nối hệ thống core.",
"In our live demo, we will walk through three clear moments. First, a read-only query asking why a Netflix payment failed. Second, a write action to enable online payments, showing the approval sheet and biometric passkey. Third, a blocked request attempting to transfer funds. Our production roadmap focuses on three steps: binding sessions to real customer identities, implementing multi-customer isolation, and integrating with GoTyme's core banking APIs.")

# Save presentation
out_pptx = os.path.join(os.path.dirname(os.path.abspath(__file__)), "Voges_Presentation.pptx")
prs.save(out_pptx)
print("Saved premium presentation:", out_pptx)
print("Total Slides:", len(prs.slides._sldIdLst))

# Write presentation notes
out_notes = os.path.join(os.path.dirname(os.path.abspath(__file__)), "Voges_Presentation_Notes.md")
notes_content = """# Voges — Premium 6-Slide Pitch Presentation Notes

Companion notes for `Voges_Presentation.pptx` (6 slides), matching the **Agentic AI Build Week 2026 Pitching Playbook** guidelines.

For each slide, you will find:
- **VN Ghi chú (VN Notes):** Thuyết minh bằng tiếng Việt để chuẩn bị tinh thần và nắm rõ mục tiêu của slide.
- **EN Speaking Script:** Lời thoại nói bằng tiếng Anh, tối ưu hóa để diễn thuyết ngắn gọn, mạch lạc trong khoảng 30s-45s/slide (tổng cộng 4.5 phút cho toàn bộ pitching).
- **Demo / Interaction:** Hướng dẫn demo và hành động trực quan tại slide đó.

---

## Slide 1 — Team & Promise (Title Slide - Dark Background)

**VN Ghi chú:**
Mở đầu ấn tượng và dứt khoát. Giới thiệu tên dự án Voges và định vị sản phẩm: một "Voice-First AI Financial Concierge" dành riêng cho GoTyme Bank. Hãy nhấn mạnh thông điệp cốt lõi quyết định kiến trúc: Giọng nói là giao diện chính để tương tác, còn màn hình điện thoại đóng vai trò là lớp bảo vệ an toàn để xác nhận.

**EN Speaking Script:**
"Good afternoon, judges. We are the team behind Voges, a voice-first AI financial concierge built specifically for GoTyme Bank. Our core architectural promise is simple: voice is the primary interface, and the screen is the safety layer. We eliminate the friction of complex banking menus while keeping your transactions completely secure. Let me show you why this is the future of mobile banking."

---

## Slide 2 — Problem Insight (Light Background)

**VN Ghi chú:**
Trình bày nỗi đau thực tế một cách sắc bén (insight). 
1. Giao diện app ngân hàng hiện tại quá nhiều menu lồng nhau (4+ tầng) khiến việc thực hiện các tính năng khẩn cấp (như khóa thẻ, mở thanh toán trực tuyến) cực kỳ chậm chạp và khó tìm.
2. Việc đưa AI vào để thực hiện giao dịch tự động bằng giọng nói lại vấp phải bài toán bảo mật: AI thông thường chỉ biết trả lời FAQ (vô dụng), còn nếu cấp quyền ghi trực tiếp cho AI (writes) thì rủi ro từ prompt injection và ảo giác (hallucination) sẽ dẫn đến thất thoát tiền của khách hàng.

**EN Speaking Script:**
"Modern banking apps are powerful but suffer from buried navigation. Crucial security settings like freezing a card or toggling online payments are hidden under 4 or 5 layers of sub-menus, causing high user anxiety and support overhead. While natural voice control is the obvious solution, financial AI faces a security dilemma: standard chatbots are FAQ-only and useless, while AI agents with direct write access are a massive liability due to hallucinations."

---

## Slide 3 — Agentic Workflow (Light Background)

**VN Ghi chú:**
Giải thích quy trình hoạt động của Voges (Goal -> Plan -> Tools -> Act -> Verify). Điểm quan trọng nhất là tính chất "Zero-Trust": AI (GPT Realtime) chỉ có quyền "Đề xuất" (Propose) hành động, còn việc "Quyết định" và "Thực thi" (Decide & Execute) hoàn toàn nằm ở phía máy chủ (Server-side Policy Engine) kết hợp xác thực sinh trắc học phần cứng (WebAuthn Passkeys).

**EN Speaking Script:**
"To solve this, we built a Zero-Trust agentic workflow. The user states their Goal naturally over a WebRTC voice stream. GPT Realtime parses the intent and proposes a structured tool call. Our server-side Tool Router intercepts it immediately. If it is a sensitive write action, our deterministic Policy Engine triggers an expiring Pending Action on screen. The user must explicitly confirm it using a native hardware Passkey, after which the server executes the secure D1 transaction."

---

## Slide 4 — Why It Wins (Dark Background)

**VN Ghi chú:**
Nêu bật 3 trụ cột kỹ thuật giúp Voges giành chiến thắng:
1. **Policy-Controlled:** Bộ lọc quy tắc backend tất định, ngăn chặn hoàn toàn việc AI bị thao túng để thực hiện chuyển khoản trái phép hoặc đổi thông tin danh tính.
2. **Passkey-Secured:** Xác thực WebAuthn chuẩn hóa. Hệ thống chỉ nhận chữ ký mã hóa từ chip bảo mật phần cứng, Voges không bao giờ tiếp cận hay lưu trữ dữ liệu vân tay/khuôn mặt của người dùng.
3. **Audit-Backed:** Mọi thao tác từ hội thoại, đề xuất tool, đánh giá policy cho tới kết quả thực thi đều được ghi lại vĩnh viễn vào D1 logs chống sửa xóa (append-only audit trail).

**EN Speaking Script:**
"Voges wins because of three strict design pillars. First, it is Policy-Controlled: deterministic backend rules override the LLM. Hallucinations cannot compromise security. Second, it is Passkey-Secured: we verify identity cryptographically using WebAuthn. Fingerprints and face data never touch our servers. Third, it is Audit-Backed: an append-only audit trail records every intent, block, and state change for total auditability."

---

## Slide 5 — Evidence + Impact (Light Background)

**VN Ghi chú:**
Đưa ra số liệu chứng minh dự án hoạt động thực tế. Voges đã vượt qua 34 kịch bản kiểm thử tự động (tests) về bảo mật, chặn đứng 100% các cuộc tấn công lừa đảo (Scam Risk Advisor). Về mặt hiệu quả kinh doanh, việc điều khiển bằng giọng nói giúp giảm tới 90% thời gian thực hiện thao tác so với tìm kiếm thủ công trên app. Toàn bộ dữ liệu demo đều được ánh xạ dựa trên các dòng sản phẩm thật của GoTyme như tài khoản Everyday Account, tài khoản tiết kiệm Go Save (lãi suất 5%), và tài khoản mua bán Vàng (PAX Gold).

**EN Speaking Script:**
"To validate credibility, we built a test suite with 34 automated scenarios covering scam risk, frozen account checks, and international payments. All 34 tests pass successfully. In terms of impact, Voges reduces settings navigation time by 90%—completing complex toggles in 5 seconds instead of 1 minute. We also seeded our database with GoTyme's actual product lines, including Everyday Accounts, high-yield Go Save goals, and tokenized PAX Gold accounts."

---

## Slide 6 — Demo + Close (Dark Background)

**VN Ghi chú:**
Chốt lại bài pitching bằng kịch bản Demo 3 bước (Xanh - Vàng - Đỏ) và Lộ trình phát triển sản xuất trung thực (Roadmap). 
- *Bước 1 (Xanh - Chỉ đọc):* Hỏi lý do Netflix bị decline.
- *Bước 2 (Vàng - Hành động nhạy cảm):* Yêu cầu bật thanh toán trực tuyến (Bật approval sheet + quét Passkey).
- *Bước 3 (Đỏ - Nguy hiểm):* Yêu cầu chuyển hết tiền và bỏ qua xác thực (Bị chặn hoàn toàn).
Lộ trình sắp tới tập trung vào bảo mật định danh động và tích hợp Core API.

**EN Speaking Script:**
"For our live demo, we will show three moments. First, a safe query explaining why a Netflix payment failed. Second, a sensitive toggle to enable online payments using a real phone passkey. Third, a dangerous request to bypass security that is blocked. Our roadmap moves from this single-client demo to core banking API integration. Talk to Voges at master.voges.pages.dev. Thank you."
"""
with open(out_notes, "w", encoding="utf-8") as f:
    f.write(notes_content.strip())
print("Saved presentation notes:", out_notes)

