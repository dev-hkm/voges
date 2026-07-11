# -*- coding: utf-8 -*-
"""
Voges hackathon presentation generator (python-pptx).

Builds a premium, minimal, mostly black/white 16:9 deck describing the Voges
"Voice-first AI Financial Concierge" product. Content is derived strictly from
the actual source code in this workspace. No app source code is modified.

Outputs:
  - Voges_Presentation.pptx
  - Voges_Presentation_Notes.md  (VN explanation + EN speaking script + demo notes)
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ----------------------------------------------------------------------------
# Design system
# ----------------------------------------------------------------------------
INK       = RGBColor(0x0A, 0x0A, 0x0A)   # near-black text
PAPER     = RGBColor(0xFF, 0xFF, 0xFF)   # white background
GRAY      = RGBColor(0x6B, 0x72, 0x80)   # secondary text
LIGHT     = RGBColor(0xE5, 0xE7, 0xEB)   # hairlines / dividers
PANEL     = RGBColor(0xF6, 0xF7, 0xF9)   # subtle panel fill
BLACKBG   = RGBColor(0x0B, 0x0B, 0x0D)   # title / dark slides

SAFE      = RGBColor(0x10, 0xB9, 0x81)   # green  - read-only / safe
APPROVE   = RGBColor(0xF5, 0x9E, 0x0B)   # amber  - approval / confirmation
BLOCK     = RGBColor(0xEF, 0x44, 0x44)   # red    - blocked / high risk

SAFE_BG   = RGBColor(0xEC, 0xFD, 0xF5)
APPROVE_BG= RGBColor(0xFF, 0xFB, 0xEB)
BLOCK_BG  = RGBColor(0xFE, 0xF2, 0xF2)

FONT      = "Segoe UI"
FONT_LT   = "Segoe UI Light"
FONT_SB   = "Segoe UI Semibold"

EMU_W = Inches(13.333)
EMU_H = Inches(7.5)

prs = Presentation()
prs.slide_width = EMU_W
prs.slide_height = EMU_H
BLANK = prs.slide_layouts[6]


# ----------------------------------------------------------------------------
# Low-level helpers
# ----------------------------------------------------------------------------
def slide():
    return prs.slides.add_slide(BLANK)


def bg(s, color):
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = color


def rect(s, x, y, w, h, fill=None, line=None, line_w=0.75, shape=MSO_SHAPE.RECTANGLE):
    sp = s.shapes.add_shape(shape, x, y, w, h)
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    return sp


def txt(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
        space_after=6, line_spacing=1.0):
    """runs: list of paragraphs; each paragraph is list of (text, size, color, bold, font)."""
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        p.line_spacing = line_spacing
        for (t, size, color, bold, fname) in para:
            r = p.add_run()
            r.text = t
            r.font.size = Pt(size)
            r.font.color.rgb = color
            r.font.bold = bold
            r.font.name = fname
    return tb


def one(t, size, color, bold=False, fname=FONT):
    return [(t, size, color, bold, fname)]


def kicker(s, text, color=GRAY, x=Inches(0.9), y=Inches(0.62)):
    txt(s, x, y, Inches(11), Inches(0.4),
        [[(text.upper(), 13, color, True, FONT_SB)]])


def headline(s, text, y=Inches(1.02), color=INK, size=34, x=Inches(0.9), w=Inches(11.5)):
    txt(s, x, y, w, Inches(1.1), [[(text, size, color, True, FONT)]])


def accent_line(s, x=Inches(0.92), y=Inches(0.95), w=Inches(0.55), color=INK):
    rect(s, x, y, w, Pt(3), fill=color)


def footer(s, idx):
    txt(s, Inches(0.9), Inches(7.02), Inches(6), Inches(0.3),
        [[("VOGES", 9, GRAY, True, FONT_SB), ("  ·  Voice-first AI Financial Concierge", 9, LIGHT, False, FONT)]])
    txt(s, Inches(11.4), Inches(7.02), Inches(1.0), Inches(0.3),
        [[(f"{idx:02d}", 9, GRAY, True, FONT_SB)]], align=PP_ALIGN.RIGHT)


def bullet(s, x, y, w, items, gap=0.62, size=15, dot=INK, tcolor=INK):
    """items: list of (bold_lead, rest) or plain string."""
    cy = y
    for it in items:
        rect(s, x, cy + Inches(0.09), Inches(0.12), Inches(0.12), fill=dot, shape=MSO_SHAPE.OVAL)
        if isinstance(it, tuple):
            lead, rest = it
            runs = [[(lead, size, tcolor, True, FONT_SB), (rest, size, GRAY, False, FONT)]]
        else:
            runs = [[(it, size, GRAY, False, FONT)]]
        txt(s, x + Inches(0.32), cy, w - Inches(0.32), Inches(0.7), runs, line_spacing=1.02)
        cy += Inches(gap)
    return cy


def notes(s, vn, en):
    tf = s.notes_slide.notes_text_frame
    tf.text = "VN — Ghi chu thuyet trinh:\n" + vn + "\n\nEN — Speaking script:\n" + en


def chip(s, x, y, w, label, color, bgc):
    c = rect(s, x, y, w, Inches(0.42), fill=bgc, line=color, line_w=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    try:
        c.adjustments[0] = 0.5
    except Exception:
        pass
    txt(s, x, y + Inches(0.045), w, Inches(0.34), [[(label, 11.5, color, True, FONT_SB)]],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# ============================================================================
# SLIDE 1 — TITLE
# ============================================================================
s = slide(); bg(s, BLACKBG)
# thin accent triad top
rect(s, Inches(0.9), Inches(0.9), Inches(0.9), Pt(4), fill=SAFE)
rect(s, Inches(1.9), Inches(0.9), Inches(0.9), Pt(4), fill=APPROVE)
rect(s, Inches(2.9), Inches(0.9), Inches(0.9), Pt(4), fill=BLOCK)

txt(s, Inches(0.9), Inches(2.35), Inches(11.5), Inches(1.6),
    [[("Voges", 84, PAPER, True, FONT)]])
txt(s, Inches(0.95), Inches(3.75), Inches(11.5), Inches(0.7),
    [[("Voice-first AI Financial Concierge", 26, RGBColor(0xC9,0xCE,0xD6), False, FONT_LT)]])

rect(s, Inches(0.95), Inches(4.62), Inches(0.55), Pt(3), fill=PAPER)
txt(s, Inches(0.95), Inches(4.85), Inches(11), Inches(0.7),
    [[("Voice is the primary interface. The screen is the safety layer.", 18, PAPER, True, FONT_SB)]])

txt(s, Inches(0.95), Inches(6.5), Inches(11), Inches(0.4),
    [[("GPT Realtime  ·  WebRTC  ·  Cloudflare Pages + D1  ·  WebAuthn Passkeys", 12.5, GRAY, False, FONT)]])
notes(s,
"Mo dau ngan gon. Voges la tro ly ngan hang dieu khien bang giong noi. Cau chot: giong noi la giao dien chinh, con man hinh la lop an toan. Nhan manh day khong phai chatbot thong thuong.",
"This is Voges, a voice-first AI financial concierge. The one line to remember: voice is the primary interface, and the screen is the safety layer. Let me show you why that matters for banking.")

# ============================================================================
# SLIDE 2 — PROBLEM
# ============================================================================
s = slide(); bg(s, PAPER)
kicker(s, "The Problem"); accent_line(s, color=BLOCK)
headline(s, "Banking apps are powerful but hard to navigate")
bullet(s, Inches(0.95), Inches(2.15), Inches(11.4), [
    ("Menu overload.  ", "Card controls, KYC, limits and statements are buried across deep settings screens."),
    ("Users get stuck.  ", "People often do not know where a setting lives, so simple tasks stall."),
    ("Chatbots only talk.  ", "Traditional assistants answer FAQs but cannot safely coordinate real actions."),
    ("Safety is the hard part.  ", "Financial AI has to be genuinely useful without ever becoming dangerous."),
], gap=0.78, size=16)
footer(s, 2)
notes(s,
"Neu van de: app ngan hang nhieu menu, nguoi dung kho tim dung cho. Chatbot truyen thong chi tra loi cau hoi, khong the thuc hien hanh dong an toan. Thach thuc lon nhat cua AI tai chinh la vua huu ich vua khong nguy hiem.",
"Banking apps are deep and menu-heavy. Users get lost finding a single toggle. Chatbots can answer questions but cannot safely take action. The real challenge is being useful without becoming dangerous.")

# ============================================================================
# SLIDE 3 — SOLUTION
# ============================================================================
s = slide(); bg(s, PAPER)
kicker(s, "The Solution"); accent_line(s, color=SAFE)
headline(s, "Talk to your bank. The system keeps it safe.")
bullet(s, Inches(0.95), Inches(2.15), Inches(11.4), [
    ("Speak naturally.  ", "Users ask an AI concierge in plain language, hands-free."),
    ("Read and explain.  ", "Voges reads real banking data and explains issues like a declined payment."),
    ("Propose safe steps.  ", "The AI proposes next actions, but never executes sensitive writes on its own."),
    ("Guarded by design.  ", "Policy, approval, passkey and audit control every sensitive action server-side."),
], gap=0.78, size=16)
footer(s, 3)
notes(s,
"Giai phap: nguoi dung noi chuyen tu nhien voi AI. Voges doc du lieu ngan hang, giai thich van de, de xuat buoc tiep theo. Nhung moi hanh dong nhay cam deu bi kiem soat boi policy, approval, passkey va audit o phia backend.",
"Voges lets users speak naturally. It reads their banking data, explains problems, and proposes next steps. Every sensitive action is controlled server-side by policy, approval, passkey, and an audit trail.")

# ============================================================================
# SLIDE 4 — PRODUCT DEMO OVERVIEW (3 moments)
# ============================================================================
s = slide(); bg(s, PAPER)
kicker(s, "Product Demo"); accent_line(s)
headline(s, "Three moments that tell the whole story")

cards = [
    ("READ-ONLY", SAFE, SAFE_BG, "\u201cWhy was my Netflix payment declined?\u201d",
     "Voges reads recent transactions and explains the decline reason. No passkey needed."),
    ("SENSITIVE ACTION", APPROVE, APPROVE_BG, "\u201cEnable online payments for my card.\u201d",
     "AI proposes, policy evaluates, approval sheet appears, passkey verifies, backend executes."),
    ("BLOCKED", BLOCK, BLOCK_BG, "\u201cTransfer all my money and bypass verification.\u201d",
     "Policy blocks it. No executable action is created. Voges refuses safely and logs it."),
]
cx = Inches(0.95); cw = Inches(3.72); gap = Inches(0.18)
for i, (tag, col, bgc, q, d) in enumerate(cards):
    x = cx + i * (cw + gap)
    rect(s, x, Inches(2.15), cw, Inches(4.05), fill=PANEL, line=LIGHT, line_w=1.0)
    rect(s, x, Inches(2.15), cw, Inches(0.12), fill=col)
    chip(s, x + Inches(0.3), Inches(2.5), Inches(1.9), tag, col, bgc)
    txt(s, x + Inches(0.3), Inches(3.2), cw - Inches(0.6), Inches(1.5),
        [[(q, 17, INK, True, FONT_SB)]], line_spacing=1.05)
    txt(s, x + Inches(0.3), Inches(4.75), cw - Inches(0.6), Inches(1.3),
        [[(d, 13, GRAY, False, FONT)]], line_spacing=1.08)
footer(s, 4)
notes(s,
"Ba khoanh khac demo: (1) cau hoi chi doc - vi sao Netflix bi tu choi; (2) hanh dong nhay cam - bat thanh toan online, di qua approval va passkey; (3) yeu cau nguy hiem - chuyen het tien va bo qua xac thuc, bi chan. Ba mau xanh/vang/do se xuat hien xuyen suot bai.",
"Three demo moments. A safe read-only question. A sensitive action that goes through approval and passkey. And a dangerous request that gets blocked. Watch the green, amber, and red pattern throughout.")

# ============================================================================
# SLIDE 5 — SYSTEM ARCHITECTURE (pipeline)
# ============================================================================
s = slide(); bg(s, PAPER)
kicker(s, "System Architecture"); accent_line(s)
headline(s, "One pipeline from voice to spoken result")

steps = [
    ("Voice", INK), ("GPT Realtime", INK), ("Tool Router", INK),
    ("Policy Engine", APPROVE), ("Approval Gate", APPROVE), ("WebAuthn", APPROVE),
    ("Execute Tool", SAFE), ("Audit Log", INK), ("Voice Response", INK),
]
# two rows of boxes with connectors
bw = Inches(2.55); bh = Inches(0.82)
xs0 = Inches(0.95); ys0 = Inches(2.35); hgap = Inches(0.42); vgap = Inches(1.15)
per_row = 3
positions = []
for i in range(len(steps)):
    row = i // per_row
    col = i % per_row
    # serpentine direction
    if row % 2 == 1:
        col = per_row - 1 - col
    x = xs0 + col * (bw + hgap)
    y = ys0 + row * (bh + vgap)
    positions.append((x, y))

for i, (label, col) in enumerate(steps):
    x, y = positions[i]
    b = rect(s, x, y, bw, bh, fill=PAPER, line=col, line_w=1.5, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    try: b.adjustments[0] = 0.12
    except Exception: pass
    txt(s, x, y + Inches(0.02), bw, bh, [[(label, 15, col, True, FONT_SB)]],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, x, y + Inches(0.02), bw, bh, None) if False else None

# arrows between consecutive
def arrow(s, p1, p2):
    x1, y1 = p1; x2, y2 = p2
    ln = s.shapes.add_connector(2, 0, 0, 0, 0)  # straight
    # decide anchor points
    if abs(int(y1) - int(y2)) < 10:  # same row
        if x2 > x1:
            sx = x1 + bw; sy = y1 + bh/2; ex = x2; ey = y2 + bh/2
        else:
            sx = x1; sy = y1 + bh/2; ex = x2 + bw; ey = y2 + bh/2
    else:  # drop down
        sx = x1 + bw/2; sy = y1 + bh; ex = x2 + bw/2; ey = y2
    ln.begin_x, ln.begin_y, ln.end_x, ln.end_y = int(sx), int(sy), int(ex), int(ey)
    ln.line.color.rgb = GRAY
    ln.line.width = Pt(1.5)
    tail = ln.line._get_or_add_ln()
    end = tail.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'})
    tail.append(end)

for i in range(len(steps) - 1):
    arrow(s, positions[i], positions[i + 1])

# legend
ly = Inches(6.4)
for lbl, col, bgc, dx in [("Safe / read-only", SAFE, SAFE_BG, 0.0),
                          ("Approval + passkey", APPROVE, APPROVE_BG, 3.1),
                          ("Blocked / high risk", BLOCK, BLOCK_BG, 6.4)]:
    rect(s, Inches(0.95 + dx), ly + Inches(0.03), Inches(0.28), Inches(0.28), fill=col, shape=MSO_SHAPE.OVAL)
    txt(s, Inches(1.32 + dx), ly, Inches(3), Inches(0.34), [[(lbl, 12.5, GRAY, False, FONT)]], anchor=MSO_ANCHOR.MIDDLE)
footer(s, 5)
notes(s,
"Kien truc mot duong ong lien mach. Giong noi vao GPT Realtime, model goi tool, tool router chay backend. Policy engine quyet dinh, cong approval, WebAuthn xac thuc, thuc thi tool, ghi audit, roi tra ket qua bang giong noi. Mau vang la buoc kiem soat, mau xanh la thuc thi that.",
"This is the full pipeline. Voice goes into GPT Realtime, the model calls a tool, the router hits the backend. The policy engine decides, the approval gate and WebAuthn verify, the tool executes, everything is audited, and the answer comes back as voice.")

# ============================================================================
# SLIDE 6 — VOICE + TOOL CALLING
# ============================================================================
s = slide(); bg(s, PAPER)
kicker(s, "Voice + Tool Calling"); accent_line(s, color=SAFE)
headline(s, "Real-time voice that can read real data")

# left column: how it works
bullet(s, Inches(0.95), Inches(2.2), Inches(6.0), [
    ("GPT Realtime.  ", "Model gpt-realtime-2.1 drives the conversation."),
    ("WebRTC.  ", "The browser streams audio directly over a peer connection."),
    ("Tool calling.  ", "The model calls typed banking tools when it needs data."),
    ("Cloudflare D1.  ", "Tools read real demo data, never invented numbers."),
], gap=0.82, size=15.5)

# right column: tool list panel
px = Inches(7.25); pw = Inches(5.1)
rect(s, px, Inches(2.1), pw, Inches(4.2), fill=PANEL, line=LIGHT, line_w=1.0)
txt(s, px + Inches(0.35), Inches(2.32), pw - Inches(0.7), Inches(0.4),
    [[("READ-ONLY BANKING TOOLS", 12, GRAY, True, FONT_SB)]])
tools = ["getCustomerProfile", "getRecentTransactions", "getCardStatus",
         "getKycStatus", "explainDeclineReason", "generateFundingInstruction"]
ty = Inches(2.85)
for t in tools:
    rect(s, px + Inches(0.35), ty + Inches(0.11), Inches(0.1), Inches(0.1), fill=SAFE, shape=MSO_SHAPE.OVAL)
    txt(s, px + Inches(0.6), ty, pw - Inches(1.0), Inches(0.4),
        [[(t + "( )", 14.5, INK, False, "Consolas")]])
    ty += Inches(0.52)
footer(s, 6)
notes(s,
"Phan voice va tool calling. Voges dung GPT Realtime model gpt-realtime-2.1, ket noi bang WebRTC truc tiep tu trinh duyet. Model goi cac tool ngan hang co dinh kieu de doc du lieu that tu Cloudflare D1. Day la cac tool chi doc, khong bia so lieu.",
"Voges uses GPT Realtime, model gpt-realtime-2.1, connected over WebRTC straight from the browser. The model calls typed banking tools to read real data from Cloudflare D1. These read tools never invent numbers.")

# ============================================================================
# SLIDE 7 — SAFETY LAYER
# ============================================================================
s = slide(); bg(s, PAPER)
kicker(s, "Safety Layer"); accent_line(s, color=APPROVE)
headline(s, "The AI proposes. The backend decides.")

# principle strip
rect(s, Inches(0.95), Inches(2.05), Inches(11.45), Inches(0.9), fill=PANEL, line=LIGHT, line_w=1.0)
txt(s, Inches(1.25), Inches(2.05), Inches(11), Inches(0.9),
    [[("A deterministic policy engine is authoritative. Frontend state is presentation only. ",
       15, INK, True, FONT_SB),
      ("Identity, permission, risk and execution are all checked server-side.", 15, GRAY, False, FONT)]],
    anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.05)

# four decision chips
dy = Inches(3.35)
decisions = [("ALLOW", SAFE, SAFE_BG, "Read-only data"),
             ("CONFIRM", APPROVE, APPROVE_BG, "Requires approval"),
             ("PASSKEY", APPROVE, APPROVE_BG, "Requires WebAuthn"),
             ("BLOCK", BLOCK, BLOCK_BG, "Refused outright")]
dw = Inches(2.72)
for i, (lab, col, bgc, sub) in enumerate(decisions):
    x = Inches(0.95) + i * (dw + Inches(0.14))
    rect(s, x, dy, dw, Inches(1.0), fill=bgc, line=col, line_w=1.2, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(s, x, dy + Inches(0.14), dw, Inches(0.4), [[(lab, 16, col, True, FONT_SB)]], align=PP_ALIGN.CENTER)
    txt(s, x, dy + Inches(0.56), dw, Inches(0.35), [[(sub, 12, GRAY, False, FONT)]], align=PP_ALIGN.CENTER)

txt(s, Inches(0.95), Inches(4.75), Inches(11), Inches(0.4),
    [[("Always blocked, by policy:", 14, INK, True, FONT_SB)]])
blocked = ["Transfer money", "Bypass verification", "Reveal CVV / OTP / full card number",
           "Change KYC identity", "Investment advice"]
bxx = Inches(0.95); byy = Inches(5.25)
for b in blocked:
    w = Inches(0.28 + 0.105 * len(b))
    rect(s, bxx, byy, w, Inches(0.5), fill=BLOCK_BG, line=BLOCK, line_w=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(s, bxx, byy, w, Inches(0.5), [[(b, 12.5, BLOCK, True, FONT_SB)]],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    bxx += w + Inches(0.16)
footer(s, 7)
notes(s,
"Lop an toan la trai tim cua Voges. AI chi de xuat, backend moi quyet dinh. Policy engine tat dinh va co tinh phan quyet: allow, confirm, passkey, hoac block. Trang thai frontend chi de hien thi. Nhung thu luon bi chan: chuyen tien, bo qua xac thuc, lo CVV/OTP, doi danh tinh KYC, tu van dau tu.",
"This is the heart of Voges. The AI only proposes; the backend decides. The policy engine is deterministic: allow, confirm, passkey, or block. Money transfers, verification bypass, revealing secrets, KYC identity changes, and investment advice are always blocked.")

# ============================================================================
# SLIDE 8 — APPROVAL + WEBAUTHN
# ============================================================================
s = slide(); bg(s, PAPER)
kicker(s, "Approval + WebAuthn"); accent_line(s, color=APPROVE)
headline(s, "A real passkey, not a simulated tap")

# left: steps
bullet(s, Inches(0.95), Inches(2.2), Inches(6.1), [
    ("Pending action.  ", "Backend creates an expiring proposal, not an execution."),
    ("Approval sheet.  ", "UI shows current state, new state, risk level and policy reason."),
    ("User confirms.  ", "The person explicitly approves the change on screen."),
    ("Passkey prompt.  ", "The OS / browser runs the real WebAuthn verification."),
    ("Server verifies.  ", "Backend validates the cryptographic assertion, then executes."),
], gap=0.72, size=15)

# right: privacy panel
px = Inches(7.35); pw = Inches(5.0)
rect(s, px, Inches(2.1), pw, Inches(3.05), fill=BLACKBG)
txt(s, px + Inches(0.4), Inches(2.4), pw - Inches(0.8), Inches(0.5),
    [[("Biometric privacy", 18, PAPER, True, FONT_SB)]])
txt(s, px + Inches(0.4), Inches(3.05), pw - Inches(0.8), Inches(2.0),
    [[("Voges never receives fingerprint or face data.", 15, PAPER, True, FONT_SB)],
     [("The operating system performs local user verification and returns only a signed cryptographic assertion.",
       14, RGBColor(0xC9,0xCE,0xD6), False, FONT)],
     [("No fake biometric. No simulated success.", 14, APPROVE, True, FONT_SB)]],
    line_spacing=1.12, space_after=10)
chip(s, px, Inches(5.45), Inches(2.4), "@simplewebauthn/server", INK, PANEL)
chip(s, px + Inches(2.6), Inches(5.45), Inches(2.4), "@simplewebauthn/browser", INK, PANEL)
footer(s, 8)
notes(s,
"Voi hanh dong nhay cam, backend tao pending action het han. UI hien approval sheet: trang thai hien tai, trang thai moi, muc rui ro, ly do policy. Nguoi dung xac nhan, roi he dieu hanh bat passkey that. Backend xac minh chu ky WebAuthn roi moi thuc thi. Voges khong nhan du lieu van tay hay khuon mat - chi nhan mot chu ky mat ma.",
"For sensitive actions the backend creates an expiring pending action. The approval sheet shows the current state, the new state, the risk, and the policy reason. The user confirms, the OS runs a real passkey, and the server verifies the cryptographic assertion before executing. Voges never sees your fingerprint or face, only a signed assertion.")

# ============================================================================
# SLIDE 9 — AUDIT TRAIL
# ============================================================================
s = slide(); bg(s, PAPER)
kicker(s, "Audit Trail"); accent_line(s)
headline(s, "Every step is written down")
bullet(s, Inches(0.95), Inches(2.2), Inches(5.9), [
    ("Append-only log.  ", "Every tool call and action is recorded."),
    ("Full context.  ", "Policy result, pending action, confirmation, passkey status, execution result."),
    ("Built for trust.  ", "Supports compliance, debugging, and after-the-fact review."),
], gap=0.9, size=15.5)

# right: timeline
px = Inches(7.2); pw = Inches(5.15)
rect(s, px, Inches(2.1), pw, Inches(4.2), fill=PANEL, line=LIGHT, line_w=1.0)
events = [("tool_read_completed", SAFE), ("action_proposed", INK),
          ("policy_evaluated", APPROVE), ("pending_action_created", APPROVE),
          ("webauthn_verified", APPROVE), ("action_executed", SAFE),
          ("policy_blocked", BLOCK)]
ey = Inches(2.45)
rect(s, px + Inches(0.62), Inches(2.6), Pt(2), Inches(3.35), fill=LIGHT)
for name, col in events:
    rect(s, px + Inches(0.5), ey + Inches(0.06), Inches(0.24), Inches(0.24), fill=col, shape=MSO_SHAPE.OVAL)
    txt(s, px + Inches(0.95), ey, pw - Inches(1.2), Inches(0.4),
        [[(name, 14, INK, False, "Consolas")]], anchor=MSO_ANCHOR.MIDDLE)
    ey += Inches(0.5)
footer(s, 9)
notes(s,
"Audit trail: moi tool call va moi action deu duoc ghi vao log chi them. Log gom ket qua policy, pending action, xac nhan, trang thai WebAuthn, ket qua thuc thi. Ke ca yeu cau bi chan cung duoc ghi lai. Day la nen tang cho tuan thu, go loi va niem tin.",
"Everything is written to an append-only audit log: the policy result, the pending action, the confirmation, the passkey status, and the execution result. Even a blocked request is logged. This is what makes the system trustworthy and auditable.")

# ============================================================================
# SLIDE 10 — DEMO SCENARIOS (6)
# ============================================================================
s = slide(); bg(s, PAPER)
kicker(s, "Demo Scenarios"); accent_line(s)
headline(s, "Six scenarios we can run live")
scen = [
    ("Card declined", "Explain why a payment failed", SAFE, SAFE_BG),
    ("Recent transactions", "Show a transaction summary", SAFE, SAFE_BG),
    ("KYC status", "Read verification state", SAFE, SAFE_BG),
    ("Funding account", "Return funding instructions", SAFE, SAFE_BG),
    ("Enable online payments", "Approval + passkey + execute", APPROVE, APPROVE_BG),
    ("Dangerous request", "Policy blocks and refuses", BLOCK, BLOCK_BG),
]
cw = Inches(3.72); ch = Inches(1.7); gx = Inches(0.18); gy = Inches(0.22)
x0 = Inches(0.95); y0 = Inches(2.2)
for i, (t, d, col, bgc) in enumerate(scen):
    r = i // 3; c = i % 3
    x = x0 + c * (cw + gx); y = y0 + r * (ch + gy)
    rect(s, x, y, cw, ch, fill=PAPER, line=LIGHT, line_w=1.0)
    rect(s, x, y, Inches(0.12), ch, fill=col)
    txt(s, x + Inches(0.38), y + Inches(0.24), cw - Inches(0.6), Inches(0.5),
        [[(t, 16.5, INK, True, FONT_SB)]])
    txt(s, x + Inches(0.38), y + Inches(0.82), cw - Inches(0.6), Inches(0.7),
        [[(d, 13, GRAY, False, FONT)]], line_spacing=1.05)
footer(s, 10)
notes(s,
"Sau kich ban demo: (1) giai thich giao dich bi tu choi; (2) hien giao dich gan day; (3) kiem tra KYC; (4) huong dan nap tien; (5) bat thanh toan online voi approval va passkey; (6) yeu cau nguy hiem bi chan. Bon cai dau la chi doc mau xanh, cai thu nam mau vang, cai cuoi mau do.",
"Six scenarios. Four safe read-only ones: explaining a decline, showing transactions, checking KYC, and funding guidance. One sensitive action with approval and passkey. And one dangerous request that gets blocked.")

# ============================================================================
# SLIDE 11 — WHAT MAKES VOGES DIFFERENT
# ============================================================================
s = slide(); bg(s, PAPER)
kicker(s, "Differentiation"); accent_line(s)
headline(s, "Not a chatbot. A guarded voice agent.")

colw = Inches(5.55)
# left: traditional chatbot
lx = Inches(0.95)
rect(s, lx, Inches(2.1), colw, Inches(4.2), fill=PANEL, line=LIGHT, line_w=1.0)
txt(s, lx + Inches(0.4), Inches(2.32), colw - Inches(0.8), Inches(0.5),
    [[("Traditional chatbot", 18, GRAY, True, FONT_SB)]])
bullet(s, lx + Inches(0.4), Inches(3.0), colw - Inches(0.8), [
    "Text-first", "Answers FAQs", "Weak action safety", "Little visibility",
], gap=0.62, size=15, dot=GRAY)

# right: voges
rx = Inches(6.85)
rect(s, rx, Inches(2.1), colw, Inches(4.2), fill=BLACKBG)
txt(s, rx + Inches(0.4), Inches(2.32), colw - Inches(0.8), Inches(0.5),
    [[("Voges", 18, PAPER, True, FONT_SB)]])
vitems = ["Voice-first", "Tool-based", "Policy-controlled", "Passkey-secured",
          "Audit-backed", "Screen as safety layer"]
cy = Inches(2.98)
for it in vitems:
    rect(s, rx + Inches(0.4), cy + Inches(0.08), Inches(0.12), Inches(0.12), fill=SAFE, shape=MSO_SHAPE.OVAL)
    txt(s, rx + Inches(0.72), cy, colw - Inches(1.1), Inches(0.4),
        [[(it, 15, PAPER, True, FONT_SB)]])
    cy += Inches(0.53)
footer(s, 11)
notes(s,
"Diem khac biet. Chatbot truyen thong: uu tien text, tra loi FAQ, an toan hanh dong yeu, it minh bach. Voges: uu tien giong noi, dua tren tool, kiem soat bang policy, bao mat bang passkey, co audit, va dung man hinh nhu lop an toan. Do chinh la cau chot mo dau.",
"A traditional chatbot is text-first, answers FAQs, has weak action safety, and little visibility. Voges is voice-first, tool-based, policy-controlled, passkey-secured, audit-backed, and uses the screen as a safety layer.")

# ============================================================================
# SLIDE 12 — ROADMAP / NEXT STEPS
# ============================================================================
s = slide(); bg(s, PAPER)
kicker(s, "Roadmap"); accent_line(s)
headline(s, "Honest next steps to production")

# demo limitation note
rect(s, Inches(0.95), Inches(2.02), Inches(11.45), Inches(0.85), fill=APPROVE_BG, line=APPROVE, line_w=1.0)
txt(s, Inches(1.25), Inches(2.02), Inches(11), Inches(0.85),
    [[("Demo limitation:  ", 14.5, APPROVE, True, FONT_SB),
      ("today the app resolves a single demo customer. Real authenticated identity is the first production step.",
       14.5, INK, False, FONT)]],
    anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.05)

left = [
    ("Authenticated identity.  ", "Bind sessions to a real, logged-in customer."),
    ("Multi-customer switching.  ", "Only after identity binding is in place."),
    ("Expanded policy rules.  ", "Cover more actions and risk conditions."),
]
right = [
    ("Production observability.  ", "Metrics, tracing, and alerting."),
    ("Bank core integration.  ", "Connect to real ledgers and card systems."),
    ("Human escalation + languages.  ", "Live handoff and broader multilingual support."),
]
bullet(s, Inches(0.95), Inches(3.2), Inches(5.7), left, gap=0.92, size=15)
bullet(s, Inches(6.85), Inches(3.2), Inches(5.5), right, gap=0.92, size=15)
footer(s, 12)
notes(s,
"Lo trinh thang than. Gioi han demo: hien tai app dung mot khach hang demo co dinh, buoc dau tien la danh tinh xac thuc that. Sau do: chuyen doi nhieu khach hang sau khi da rang buoc danh tinh, mo rong luat policy, observability production, tich hop core ngan hang, luong chuyen nguoi that va da ngon ngu. Khong noi qua, chi noi dung nhung gi con thieu.",
"Let me be honest about what is next. Today the demo resolves a single customer, so real authenticated identity is the first step. After that: multi-customer switching, more policy rules, production observability, bank core integration, human escalation, and better multilingual support.")

# ----------------------------------------------------------------------------
out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "Voges_Presentation.pptx")
prs.save(out)
print("Saved:", out)
print("Slides:", len(prs.slides._sldIdLst))
